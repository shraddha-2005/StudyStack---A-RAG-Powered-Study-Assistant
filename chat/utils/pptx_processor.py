from pptx import Presentation
from pptx.util import Inches
from pathlib import Path
import re


def process_pptx(file_path, original_filename):
    """Extract text from PowerPoint presentations"""
    try:
        prs = Presentation(file_path)
        title = Path(original_filename).stem
        full_text = []
        slide_texts = []

        for slide_num, slide in enumerate(prs.slides, start=1):
            slide_content = []

            slide_title = None
            if slide.shapes.title and slide.shapes.title.text.strip():
                slide_title = slide.shapes.title.text.strip()

            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        slide_content.append(text)

            if slide_content:
                header = f"[Slide {slide_num}" + (f": {slide_title}]" if slide_title else "]")
                combined = header + "\n" + "\n".join(slide_content)
                slide_texts.append(combined)
                full_text.append(combined)

        if not full_text:
            return {'success': False, 'error': 'No text content found in presentation'}

        chunks = []
        group_size = 3
        for i in range(0, len(slide_texts), group_size):
            group = slide_texts[i:i + group_size]
            chunk_text = "\n\n".join(group)
            start_slide = i + 1
            end_slide = min(i + group_size, len(slide_texts))

            chunks.append({
                'text': chunk_text,
                'title': f"{title} (PPTX)",
                'number': f"Slides {start_slide}-{end_slide}",
                'start': start_slide,
                'end': end_slide,
            })

        return {
            'success': True,
            'title': title,
            'text': "\n\n".join(full_text),
            'chunks': chunks,
            'slide_count': len(prs.slides),
        }

    except Exception as e:
        return {'success': False, 'error': str(e)}